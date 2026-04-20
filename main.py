import subprocess
import sys

_PACKAGES = {
    "fastapi":           "fastapi==0.115.0",
    "uvicorn":           "uvicorn[standard]==0.30.6",
    "multipart":         "python-multipart==0.0.12",
    "dotenv":            "python-dotenv",
    "anthropic":         "anthropic",
    "pptx":              "python-pptx",
    "pdfplumber":        "pdfplumber",
    "docx":              "python-docx",
    "openpyxl":          "openpyxl",
    "xlrd":              "xlrd==1.2.0",
}

for _import_name, _pkg_name in _PACKAGES.items():
    try:
        __import__(_import_name)
    except ImportError:
        print(f"[자동설치] {_pkg_name} 설치 중...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", _pkg_name])

import io
import os
import re
import json
import tempfile
import anthropic

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from dotenv import load_dotenv

load_dotenv()

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

app = FastAPI(title="AI PPT Generator")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

app.mount("/static", StaticFiles(directory=os.path.join(BASE_DIR, "static")), name="static")

@app.get("/")
def root():
    return FileResponse(os.path.join(BASE_DIR, "static", "index.html"))


# ── 테마 정의 ────────────────────────────────────────────────────
THEMES = {
    "dark":      {"bg": "1E2761", "header": "0D1B5E", "title": "FFFFFF", "body": "CADCFC", "accent": "4FC3F7"},
    "corporate": {"bg": "1565C0", "header": "0D47A1", "title": "FFFFFF", "body": "E3F2FD", "accent": "FFD54F"},
    "modern":    {"bg": "4A148C", "header": "38006B", "title": "FFFFFF", "body": "EDE7F6", "accent": "CE93D8"},
    "warm":      {"bg": "B85042", "header": "8D2B1A", "title": "FFFFFF", "body": "ECE2D0", "accent": "F9E795"},
}

def hex_to_rgb(hex_str: str) -> RGBColor:
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, color_hex):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, color_hex, font_size, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color_hex)
    run.font.name = font_name
    return txBox


# ── 엑셀 구조 분석 ────────────────────────────────────────────────
def extract_excel_structure(content_bytes: bytes, filename: str) -> str:
    """
    Excel(xlsx/xlsm/xls)을 분석하여 Claude가 스토리를 파악할 수 있는
    구조화된 텍스트 요약을 생성합니다.
    - xlsx/xlsm: openpyxl로 두 번 로드(수식 + 계산값)
    - xls      : xlrd 사용
    """
    ext = filename.rsplit(".", 1)[-1].lower()

    # ── xls 구버전 ────────────────────────────────────────────────
    if ext == "xls":
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=content_bytes)
            parts = [f"[Excel 파일: {filename}]  형식: xls(구버전)"]
            for sname in wb.sheet_names():
                ws = wb.sheet_by_name(sname)
                parts.append(f"\n## 시트: '{sname}'  ({ws.nrows}행 × {ws.ncols}열)")
                if ws.nrows == 0:
                    parts.append("  (빈 시트)")
                    continue
                for r in range(min(ws.nrows, 60)):
                    row_vals = []
                    for c in range(min(ws.ncols, 25)):
                        cell = ws.cell(r, c)
                        v = cell.value
                        if v == "":
                            row_vals.append("")
                        elif cell.ctype == xlrd.XL_CELL_NUMBER:
                            row_vals.append(str(int(v)) if v == int(v) else str(round(v, 4)))
                        else:
                            row_vals.append(str(v))
                    line = " | ".join(v for v in row_vals if v)
                    if line.strip():
                        parts.append(f"  {line}")
            return "\n".join(parts)
        except ImportError:
            return f"[{filename}: xls 읽기 실패 — pip install xlrd==1.2.0 필요]"
        except Exception as e:
            return f"[xls 파싱 실패: {e}]"

    # ── xlsx / xlsm ───────────────────────────────────────────────
    try:
        import openpyxl
        keep_vba = (ext == "xlsm")
        wb_val = openpyxl.load_workbook(io.BytesIO(content_bytes), data_only=True,  keep_vba=keep_vba)
        wb_fml = openpyxl.load_workbook(io.BytesIO(content_bytes), data_only=False, keep_vba=keep_vba)
    except Exception as e:
        return f"[Excel 파싱 실패: {e}]"

    parts = [f"[Excel 파일 분석: {filename}]  형식: {ext.upper()}"]
    if ext == "xlsm":
        parts.append("(VBA 매크로 포함 — 데이터/구조만 분석)")

    # Named Ranges — 작성자가 중요하다고 이름 붙인 셀
    try:
        named = [(n, d.attr_text) for n, d in wb_val.defined_names.items()
                 if not n.startswith("_")]
        if named:
            parts.append("\n## Named Ranges (핵심 데이터 포인터)")
            for n, ref in named[:15]:
                parts.append(f"  - {n}: {ref}")
    except Exception:
        pass

    # 시트별 분석
    for sname in wb_val.sheetnames:
        ws_v = wb_val[sname]
        ws_f = wb_fml[sname]
        max_row = ws_v.max_row or 0
        max_col = ws_v.max_column or 0

        parts.append(f"\n## 시트: '{sname}'  ({max_row}행 × {max_col}열)")
        if max_row == 0 or max_col == 0:
            parts.append("  (빈 시트)")
            continue

        # 이미 있는 차트 — 작성자의 시각화 의도
        try:
            chart_descs = []
            for ch in ws_v._charts:
                ch_type  = type(ch).__name__.replace("Chart", "")
                ch_title = getattr(ch, "title", None) or "제목없음"
                chart_descs.append(f"{ch_type}({ch_title})")
            if chart_descs:
                parts.append(f"  이미 있는 차트: {', '.join(chart_descs)}")
        except Exception:
            pass

        # 조건부 서식 — 강조된 중요 영역
        try:
            cf_ranges = [str(r) for r in ws_v.conditional_formatting][:5]
            if cf_ranges:
                parts.append(f"  강조(조건부서식) 범위: {', '.join(cf_ranges)}")
        except Exception:
            pass

        # 셀 데이터 수집
        row_limit = min(max_row, 80)
        col_limit = min(max_col, 25)
        func_counter: dict = {}
        key_formulas: list = []
        grid_rows: list = []

        IMPORTANT_FUNCS = {
            "SUM", "AVERAGE", "COUNT", "COUNTA", "MAX", "MIN",
            "VLOOKUP", "HLOOKUP", "INDEX", "MATCH",
            "IF", "IFS", "SUMIF", "SUMIFS", "COUNTIF", "COUNTIFS",
            "IFERROR", "ROUND", "SUBTOTAL",
        }

        for r in range(1, row_limit + 1):
            row_display = []
            has_content = False
            for c in range(1, col_limit + 1):
                cell_v = ws_v.cell(row=r, column=c)
                cell_f = ws_f.cell(row=r, column=c)
                val = cell_v.value
                fml = cell_f.value
                is_formula = isinstance(fml, str) and fml.startswith("=")

                if is_formula:
                    # 수식 함수명 집계
                    for fn in re.findall(r'[A-Z]{2,}(?=\()', fml.upper()):
                        func_counter[fn] = func_counter.get(fn, 0) + 1
                    # 중요 수식 저장
                    if any(fn in fml.upper() for fn in IMPORTANT_FUNCS) and len(key_formulas) < 35:
                        key_formulas.append((cell_v.coordinate, fml, val))
                    display = f"{val}" if val is not None else "(수식계산중)"
                else:
                    display = str(val) if val is not None else ""

                row_display.append(display)
                if display and display not in ("", "None"):
                    has_content = True

            if has_content:
                grid_rows.append(row_display)

        # 데이터 출력
        if grid_rows:
            parts.append(f"  데이터 ({len(grid_rows)}개 행):")
            for row in grid_rows[:50]:
                line = " | ".join(v for v in row if v and v != "None")
                if line.strip():
                    parts.append(f"    {line}")
            if len(grid_rows) > 50:
                parts.append(f"    ... (이하 {len(grid_rows) - 50}행 생략)")

        # 수식 함수 사용 현황
        if func_counter:
            summary = ", ".join(
                f"{fn}({cnt}회)"
                for fn, cnt in sorted(func_counter.items(), key=lambda x: -x[1])[:10]
            )
            parts.append(f"  수식 함수 사용: {summary}")

        # 핵심 수식 목록
        if key_formulas:
            parts.append(f"  핵심 수식 ({len(key_formulas)}개):")
            for addr, fml, val in key_formulas[:20]:
                parts.append(f"    {addr}: {fml}  →  결과: {val}")

    wb_val.close()
    wb_fml.close()
    return "\n".join(parts)


# ── 파일 텍스트 추출 ──────────────────────────────────────────────
def extract_text_from_file(content_bytes: bytes, filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
                texts = []
                for page in pdf.pages[:20]:
                    t = page.extract_text()
                    if t:
                        texts.append(t)
            return "\n\n".join(texts) if texts else f"[{filename}: PDF에서 텍스트를 추출할 수 없습니다]"
        except Exception as e:
            return f"[PDF 파싱 실패: {e}]"

    elif ext == "docx":
        try:
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
            return "\n".join(parts)
        except Exception as e:
            return f"[DOCX 파싱 실패: {e}]"

    elif ext == "csv":
        try:
            import csv
            text = content_bytes.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text))
            rows = list(reader)
            return "\n".join([" | ".join(row) for row in rows[:150]])
        except Exception as e:
            return f"[CSV 파싱 실패: {e}]"

    elif ext == "json":
        try:
            text = content_bytes.decode("utf-8", errors="replace")
            data = json.loads(text)
            return json.dumps(data, ensure_ascii=False, indent=2)[:8000]
        except Exception:
            try:
                return content_bytes.decode("utf-8", errors="replace")
            except Exception as e:
                return f"[JSON 파싱 실패: {e}]"

    elif ext in ("xlsx", "xlsm", "xls"):
        return extract_excel_structure(content_bytes, filename)

    else:
        for enc in ("utf-8", "cp949", "euc-kr"):
            try:
                return content_bytes.decode(enc)
            except UnicodeDecodeError:
                continue
        return f"[파일명: {filename}] 바이너리 파일 - 파일 이름 기반으로 프레젠테이션을 생성합니다."


# ── 차트 슬라이드 ─────────────────────────────────────────────────
CHART_TYPE_MAP = {
    "bar":     XL_CHART_TYPE.BAR_CLUSTERED,
    "column":  XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":    XL_CHART_TYPE.LINE,
    "pie":     XL_CHART_TYPE.PIE,
    "area":    XL_CHART_TYPE.AREA,
}

def add_chart_slide(slide, slide_info, t):
    add_rect(slide, 0, 0, 10, 0.9, t["header"])
    add_text(slide, slide_info.get("title", ""), 0.4, 0.1, 9.2, 0.7, t["title"], 22, bold=True)

    categories = slide_info.get("categories", [])
    series_list = slide_info.get("series", [])

    if not categories or not series_list:
        add_text(slide, "차트 데이터가 없습니다", 0.5, 2.5, 9.0, 1.0, t["body"], 14, align=PP_ALIGN.CENTER)
        return

    cd = ChartData()
    cd.categories = [str(c) for c in categories]
    for s in series_list:
        values = tuple(float(v) if v is not None else 0.0 for v in s.get("values", []))
        cd.add_series(s.get("name", ""), values)

    xl_type = CHART_TYPE_MAP.get(slide_info.get("chart_type", "column").lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

    try:
        chart_shape = slide.shapes.add_chart(
            xl_type,
            Inches(0.5), Inches(1.05),
            Inches(9.0), Inches(4.3),
            cd
        )
        chart = chart_shape.chart

        try:
            chart.chart_area.format.fill.solid()
            chart.chart_area.format.fill.fore_color.rgb = hex_to_rgb(t["bg"])
        except Exception:
            pass
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = hex_to_rgb(t["bg"])
        except Exception:
            pass

        series_colors = [t["accent"], "FFFFFF", t["body"], "FF8C00", "00CED1"]
        try:
            for i, series in enumerate(chart.series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = hex_to_rgb(series_colors[i % len(series_colors)])
        except Exception:
            pass

        try:
            if xl_type == XL_CHART_TYPE.PIE or len(series_list) > 1:
                chart.has_legend = True
                chart.legend.font.size = Pt(11)
                chart.legend.font.color.rgb = hex_to_rgb(t["body"])
            else:
                chart.has_legend = False
        except Exception:
            pass

        try:
            chart.category_axis.tick_labels.font.color.rgb = hex_to_rgb(t["body"])
            chart.value_axis.tick_labels.font.color.rgb = hex_to_rgb(t["body"])
        except Exception:
            pass  # 파이 차트는 축 없음

    except Exception as e:
        add_text(slide, f"차트 생성 오류: {str(e)[:80]}", 0.5, 2.5, 9.0, 1.0, t["body"], 11)


# ── 테이블 슬라이드 ───────────────────────────────────────────────
def add_table_slide(slide, slide_info, t):
    add_rect(slide, 0, 0, 10, 0.9, t["header"])
    add_text(slide, slide_info.get("title", ""), 0.4, 0.1, 9.2, 0.7, t["title"], 22, bold=True)

    headers = slide_info.get("headers", [])
    rows = slide_info.get("rows", [])

    if not headers:
        return

    cols = len(headers)
    total_rows = len(rows) + 1

    table_shape = slide.shapes.add_table(
        total_rows, cols,
        Inches(0.4), Inches(1.1),
        Inches(9.2), Inches(4.1)
    )
    tbl = table_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(t["header"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(h)
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = hex_to_rgb(t["accent"])
        run.font.name = "Calibri"

    for i, row_data in enumerate(rows):
        for j in range(cols):
            val = row_data[j] if j < len(row_data) else ""
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(t["bg"] if i % 2 == 0 else t["header"])
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            run.font.color.rgb = hex_to_rgb(t["body"])
            run.font.name = "Calibri"


# ── 프로세스 슬라이드 ─────────────────────────────────────────────
def add_process_slide(slide, slide_info, t):
    add_rect(slide, 0, 0, 10, 0.9, t["header"])
    add_text(slide, slide_info.get("title", ""), 0.4, 0.1, 9.2, 0.7, t["title"], 22, bold=True)

    steps = slide_info.get("steps", [])[:5]
    if not steps:
        return

    n = len(steps)
    step_w = 9.0 / n
    start_x = 0.5

    for i, step in enumerate(steps):
        x = start_x + i * step_w
        box_w = step_w - 0.25

        add_rect(slide, x, 1.5, box_w, 0.65, t["accent"])
        add_text(slide, step.get("number", f"{i+1:02d}"), x, 1.5, box_w, 0.65,
                 t["bg"], 18, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, step.get("title", ""), x, 2.3, box_w, 0.55,
                 t["title"], 12, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, step.get("desc", ""), x, 2.95, box_w, 1.5,
                 t["body"], 10, align=PP_ALIGN.CENTER)

        if i < n - 1:
            add_text(slide, "▶", x + box_w + 0.02, 1.6, 0.22, 0.5,
                     t["accent"], 14, align=PP_ALIGN.CENTER)


# ── 통계/KPI 슬라이드 ─────────────────────────────────────────────
def add_stats_slide(slide, slide_info, t):
    add_rect(slide, 0, 0, 10, 0.9, t["header"])
    add_text(slide, slide_info.get("title", ""), 0.4, 0.1, 9.2, 0.7, t["title"], 22, bold=True)

    stats = slide_info.get("stats", [])[:4]
    if not stats:
        return

    n = len(stats)
    card_w = 9.0 / n
    start_x = 0.5

    for i, stat in enumerate(stats):
        x = start_x + i * card_w
        w = card_w - 0.2

        add_rect(slide, x, 1.2, w, 3.5, t["header"])
        add_rect(slide, x, 1.2, w, 0.08, t["accent"])

        add_text(slide, stat.get("value", ""), x + 0.1, 1.5, w - 0.2, 1.3,
                 t["accent"], 36, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, stat.get("label", ""), x + 0.1, 2.9, w - 0.2, 0.6,
                 t["title"], 13, bold=True, align=PP_ALIGN.CENTER)

        if stat.get("desc"):
            add_text(slide, stat["desc"], x + 0.1, 3.55, w - 0.2, 0.7,
                     t["body"], 10, align=PP_ALIGN.CENTER)


# ── PPT 빌더 ─────────────────────────────────────────────────────
def build_pptx(slides_data: dict, theme_key: str) -> str:
    t = THEMES.get(theme_key, THEMES["dark"])
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    blank_layout = prs.slide_layouts[6]

    for i, slide_info in enumerate(slides_data.get("slides", [])):
        slide = prs.slides.add_slide(blank_layout)
        stype = slide_info.get("type", "content")

        # 새 슬라이드 타입 처리
        if stype == "chart":
            add_rect(slide, 0, 0, 10, 5.625, t["bg"])
            add_chart_slide(slide, slide_info, t)
            add_text(slide, str(i + 1), 9.0, 4.9, 0.6, 0.35, t["accent"], 10, align=PP_ALIGN.RIGHT)

        elif stype == "table":
            add_rect(slide, 0, 0, 10, 5.625, t["bg"])
            add_table_slide(slide, slide_info, t)
            add_text(slide, str(i + 1), 9.0, 4.9, 0.6, 0.35, t["accent"], 10, align=PP_ALIGN.RIGHT)

        elif stype == "process":
            add_rect(slide, 0, 0, 10, 5.625, t["bg"])
            add_process_slide(slide, slide_info, t)
            add_text(slide, str(i + 1), 9.0, 4.9, 0.6, 0.35, t["accent"], 10, align=PP_ALIGN.RIGHT)

        elif stype == "stats":
            add_rect(slide, 0, 0, 10, 5.625, t["bg"])
            add_stats_slide(slide, slide_info, t)
            add_text(slide, str(i + 1), 9.0, 4.9, 0.6, 0.35, t["accent"], 10, align=PP_ALIGN.RIGHT)

        elif stype == "title":
            add_rect(slide, 0, 0, 10, 5.625, t["bg"])
            add_rect(slide, 0.4, 2.0, 0.08, 1.5, t["accent"])
            add_text(slide, slide_info.get("title", slides_data.get("title", "")),
                     0.6, 1.9, 8.8, 1.0, t["title"], 36, bold=True)
            add_text(slide, slide_info.get("subtitle", slides_data.get("subtitle", "")),
                     0.6, 3.0, 8.8, 0.7, t["body"], 16)
            add_rect(slide, 0, 4.5, 10, 0.6, t["header"])
            add_text(slide, slides_data.get("author", ""),
                     0.4, 4.55, 9.2, 0.5, t["body"], 11)

        elif stype == "closing":
            add_rect(slide, 0, 0, 10, 5.625, t["header"])
            add_text(slide, slide_info.get("title", "감사합니다"),
                     0.5, 1.4, 9.0, 1.2, t["title"], 40, bold=True, align=PP_ALIGN.CENTER)
            if slide_info.get("message"):
                add_text(slide, slide_info["message"],
                         0.5, 2.9, 9.0, 0.8, t["accent"], 18, align=PP_ALIGN.CENTER)
            if slide_info.get("contact"):
                add_text(slide, slide_info["contact"],
                         0.5, 3.9, 9.0, 0.5, t["body"], 13, align=PP_ALIGN.CENTER)

        elif stype == "two-column":
            add_rect(slide, 0, 0, 10, 5.625, t["bg"])
            add_rect(slide, 0, 0, 10, 0.9, t["header"])
            add_text(slide, slide_info.get("title", ""), 0.4, 0.1, 9.2, 0.7, t["title"], 22, bold=True)
            add_rect(slide, 4.95, 1.0, 0.08, 4.0, t["accent"])
            left = slide_info.get("left", {})
            add_text(slide, left.get("heading", ""), 0.3, 1.0, 4.4, 0.5, t["accent"], 14, bold=True)
            for j, pt in enumerate(left.get("points", [])):
                add_text(slide, "• " + pt, 0.3, 1.6 + j * 0.6, 4.4, 0.55, t["body"], 13)
            right = slide_info.get("right", {})
            add_text(slide, right.get("heading", ""), 5.2, 1.0, 4.4, 0.5, t["accent"], 14, bold=True)
            for j, pt in enumerate(right.get("points", [])):
                add_text(slide, "• " + pt, 5.2, 1.6 + j * 0.6, 4.4, 0.55, t["body"], 13)
            add_text(slide, str(i + 1), 9.0, 4.9, 0.6, 0.35, t["accent"], 10, align=PP_ALIGN.RIGHT)

        else:  # content
            add_rect(slide, 0, 0, 10, 5.625, t["bg"])
            add_rect(slide, 0, 0, 10, 0.9, t["header"])
            add_text(slide, slide_info.get("title", ""), 0.4, 0.1, 9.2, 0.7, t["title"], 22, bold=True)
            add_rect(slide, 0.4, 1.0, 1.2, 0.05, t["accent"])
            for j, bullet in enumerate(slide_info.get("bullets", [])):
                add_rect(slide, 0.4, 1.2 + j * 0.78, 0.06, 0.38, t["accent"])
                add_text(slide, bullet, 0.65, 1.2 + j * 0.78, 9.0, 0.65, t["body"], 14)
            add_text(slide, str(i + 1), 9.0, 4.9, 0.6, 0.35, t["accent"], 10, align=PP_ALIGN.RIGHT)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


# ── Claude API 호출 ───────────────────────────────────────────────
SYSTEM_PROMPT = """당신은 전문 프레젠테이션 디자이너입니다. 파일 내용을 분석하여 시각적으로 풍부하고 설득력 있는 프레젠테이션을 만들어주세요.

사용 가능한 슬라이드 타입:
- "title": 표지 슬라이드 (항상 첫 장)
- "content": 글머리 기호 목록 (일반 내용)
- "two-column": 두 열 비교 슬라이드
- "chart": 차트/그래프 (수치·통계 데이터가 있을 때 필수 사용)
- "table": 표 슬라이드 (비교표, 데이터표)
- "process": 단계별 프로세스 다이어그램 (절차·단계 설명)
- "stats": KPI/통계 카드 (핵심 수치·지표)
- "closing": 마무리 슬라이드 (항상 마지막 장)

핵심 규칙:
1. 수치 데이터나 통계가 있으면 반드시 chart 또는 stats 슬라이드를 포함
2. 단계별 프로세스·절차가 있으면 process 슬라이드 사용
3. 비교 데이터가 있으면 table 또는 two-column 사용
4. 단순 텍스트만 있어도 chart/stats/process 타입을 창의적으로 추론하여 삽입
5. content 타입만 나열하는 것을 피하고 다양한 타입을 혼합할 것

[Excel 파일 전용 분석 규칙]
Excel 파일이 입력된 경우 다음 순서로 분석하세요:
STEP 1 — 의도 파악: 수식 함수(SUM/VLOOKUP/IF 등), Named Ranges, 이미 있는 차트를 보고
         "이 Excel이 궁극적으로 보여주려는 핵심 메시지"가 무엇인지 먼저 결론을 내리세요.
STEP 2 — 핵심 데이터 선별: 중간 계산값은 버리고, 최종 집계/비율/KPI만 추출하세요.
         조건부 서식으로 강조된 범위 = 작성자가 중요하다고 표시한 것입니다.
STEP 3 — 시트 역할 구분: 여러 시트가 있으면 '입력 시트'는 무시하고
         '요약/대시보드' 성격의 시트를 중심으로 PPT를 구성하세요.
STEP 4 — 슬라이드 배치: 파악한 스토리를 chart/stats/table을 우선 활용해 시각화하세요.

순수 JSON만 응답하세요. 마크다운 코드블록 없이."""


def call_claude(file_content: str, file_name: str, slide_count: int, language: str, purpose: str) -> dict:
    api_key = os.getenv("Anthropic_API_KEY_Win001")
    if not api_key:
        raise HTTPException(status_code=500, detail="Anthropic_API_KEY_Win001가 .env에 없습니다.")

    client = anthropic.Anthropic(api_key=api_key)

    purpose_map = {
        "general":   "일반 발표",
        "business":  "비즈니스 제안서",
        "education": "교육/강의 자료",
        "research":  "연구 보고서",
    }

    is_excel = file_name.rsplit(".", 1)[-1].lower() in ("xlsx", "xlsm", "xls")
    excel_instruction = """
[Excel 분석 지시]
위 Excel 구조를 보고, 먼저 이 파일이 전달하려는 핵심 스토리를 파악하세요.
복잡한 셀 참조와 수식의 최종 결과값, Named Ranges, 이미 있는 차트를 근거로
"이 데이터가 말하고 싶은 것"을 요약한 뒤 PPT 슬라이드를 설계하세요.
중간 계산 과정은 생략하고, 청중이 바로 이해할 수 있는 최종 인사이트만 담으세요.
""" if is_excel else ""

    user_prompt = f"""파일명: {file_name}
발표 목적: {purpose_map.get(purpose, "일반 발표")}
슬라이드 수: {slide_count}장
언어: {language}
{excel_instruction}
파일 내용:
{file_content[:6000]}

아래 JSON 형식으로 응답하세요. 총 {slide_count}장 (title 1 + 내용 여러 장 + closing 1):

{{
  "title": "프레젠테이션 제목",
  "subtitle": "부제목",
  "author": "작성자 또는 출처",
  "slides": [
    {{"type":"title","title":"발표 제목","subtitle":"부제목"}},
    {{"type":"content","title":"슬라이드 제목","bullets":["핵심 내용 1","핵심 내용 2","핵심 내용 3"]}},
    {{"type":"chart","title":"차트 제목","chart_type":"column","categories":["1분기","2분기","3분기","4분기"],"series":[{{"name":"매출(억)","values":[120,145,160,180]}}]}},
    {{"type":"table","title":"비교표","headers":["항목","A안","B안","비고"],"rows":[["비용","100만","150만","A 유리"],["기간","3개월","2개월","B 유리"]]}},
    {{"type":"process","title":"프로세스","steps":[{{"number":"01","title":"기획","desc":"요구사항 수집"}},{{"number":"02","title":"설계","desc":"시스템 설계"}},{{"number":"03","title":"개발","desc":"구현 및 테스트"}},{{"number":"04","title":"배포","desc":"서비스 출시"}}]}},
    {{"type":"stats","title":"핵심 지표","stats":[{{"value":"98%","label":"고객 만족도","desc":"+3% YoY"}},{{"value":"2.4M","label":"월간 사용자","desc":"+40% YoY"}},{{"value":"150+","label":"파트너사","desc":"글로벌 확장"}}]}},
    {{"type":"two-column","title":"비교","left":{{"heading":"장점","points":["항목1","항목2"]}},"right":{{"heading":"단점","points":["항목1","항목2"]}}}},
    {{"type":"closing","title":"감사합니다","message":"핵심 메시지","contact":"Q&A"}}
  ]
}}"""

    message = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=8192,
        system=[{
            "type": "text",
            "text": SYSTEM_PROMPT,
            "cache_control": {"type": "ephemeral"}
        }],
        messages=[{"role": "user", "content": user_prompt}]
    )

    raw = message.content[0].text.strip()
    # JSON 펜스 제거 (혹시 포함된 경우)
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    raw = raw.strip()

    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        # 응답이 잘렸을 경우: closing 슬라이드를 강제 삽입하고 닫기 시도
        # 마지막 완전한 슬라이드 객체까지만 잘라내기
        last_brace = raw.rfind("},")
        if last_brace == -1:
            last_brace = raw.rfind("}")
        if last_brace != -1:
            trimmed = raw[:last_brace + 1]
            # slides 배열 닫기 + 최상위 객체 닫기
            closing = ',{"type":"closing","title":"감사합니다","message":"","contact":"Q&A"}]}'
            try:
                return json.loads(trimmed + closing)
            except json.JSONDecodeError:
                pass
        raise


# ── 엔드포인트 ────────────────────────────────────────────────────
@app.post("/generate")
async def generate(
    file: UploadFile = File(...),
    slide_count: int = Form(8),
    language: str = Form("Korean"),
    theme: str = Form("dark"),
    purpose: str = Form("general"),
):
    content_bytes = await file.read()
    file_content = extract_text_from_file(content_bytes, file.filename)

    try:
        slides_data = call_claude(file_content, file.filename, slide_count, language, purpose)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"AI 응답 파싱 오류: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    pptx_path = build_pptx(slides_data, theme)
    safe_title = slides_data.get("title", "presentation").replace(" ", "_")[:40]
    from urllib.parse import quote
    encoded_title = quote(safe_title)

    return FileResponse(
        pptx_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{encoded_title}.pptx",
        headers={
            "X-Slide-Count": str(len(slides_data.get("slides", []))),
            "X-Presentation-Title": encoded_title,
        }
    )


@app.get("/health")
def health():
    key = os.getenv("Anthropic_API_KEY_Win001", "")
    return {"status": "ok", "api_key_set": bool(key)}


if __name__ == "__main__":
    import uvicorn
    import threading
    import webbrowser
    threading.Timer(1.5, lambda: webbrowser.open("http://localhost:8000")).start()
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
